from pptx import Presentation
from pptx.parts.chart import ChartPart
from pptx.chart.data import CategoryChartData,ChartData
from pptx.parts.embeddedpackage import EmbeddedXlsxPart
from pptx.table import _Cell,Table, _Row, _Column
from pptx.dml.color import RGBColor
from pptx.oxml.xmlchemy import OxmlElement
from pptx.slide import Slide
from pptx.enum.shapes import MSO_SHAPE_TYPE
from cast_common.logger import Logger, INFO,DEBUG

from os.path import abspath,exists
from copy import deepcopy

from pandas import DataFrame

class PowerPoint():

    log = None
    ppt = None

    def __init__(self,template:str,output:str,log_level=INFO) -> None:

        if PowerPoint.log is None:
            PowerPoint.log = Logger('PowerPoint',log_level)

        self._template_name = abspath(template)
        if not exists(self._template_name):
            raise ValueError(f'Template {self._template_name} not found!')

        self._output = abspath(output)
        self.log.info(f'Generating deck {self._output} from {self._template_name}')

        self._prs = Presentation(self._template_name)
        PowerPoint.ppt = self

    def save(self):
        self._prs.save(self._output)

    def delete_run(self,run):
        r = run._r
        r.getparent().remove(r)

    def get_shape_by_name(self, name, use_slide=None, all=False):
        rslt = []

        slides = self._prs.slides
        if use_slide != None:
            slides = [use_slide] 

        for slide in slides:
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    for grp_shape in shape.shapes:
                        parts = grp_shape.name.split(':')
                        shape_name = parts[0]
                        if shape_name == name:
                            if not all:
                                return grp_shape
                            else:
                                rslt.append(grp_shape)
                else:
                    parts = shape.name.split(':')
                    shape_name = parts[0]
                    if shape_name == name:
                        if not all:
                            return shape
                        else:
                            rslt.append(shape)

        if all:
            return rslt
        else:
            return None

    """ **************************************************************************************************************
                                            Text Replace Functionality 
    ************************************************************************************************************** """
    def replace_text (self, search_str:str, repl:any, tbd_for_blanks=True,slide=None):
        if tbd_for_blanks:
            skip = False
            omit_list = ["immediate_action","other","risk_detail"]
            for s in omit_list:
                if s in search_str:
                    skip=True

            if repl is not None and (type(repl)==int or type(repl)==float):
                repl=f'{repl:,}'

            if not skip and len(str(repl)) == 0:
                repl = 'TBD'

        if slide is None:
            for s in self._prs.slides:
                self._replace_slide_text(s, search_str, repl)
        else:
            self._replace_slide_text(slide, search_str, repl)

    def replace_textbox (self, shape_name:str, repl:str, slide:int=None):
        #get all shape with the provided name
        found = False
        shapes = self.get_shape_by_name(shape_name,use_slide=slide,all=True)
        for shape in shapes:
            if not shape is None:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        self._replace_paragraph_text(paragraph,paragraph.text,repl)
                        found=True
        return found

    def _replace_slide_text (self, slide, search_str, repl):
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                for grp_shape in shape.shapes:
                    if grp_shape.has_text_frame:
                        for paragraph in grp_shape.text_frame.paragraphs:
                            if search_str in paragraph.text:
                                self._replace_paragraph_text(paragraph,search_str,repl)
            elif shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    if search_str in paragraph.text:
                        self._replace_paragraph_text(paragraph,search_str,repl)
            elif shape.has_table:
                tbl=shape.table
                row_count = len(tbl.rows)
                col_count = len(tbl.columns)
                for r in range(0,row_count):
                    for c in range(0, col_count):
                        cell = tbl.cell(r,c)
                        for paragraph in cell.text_frame.paragraphs:
                            self._replace_paragraph_text(paragraph,search_str,repl)
        pass

    def _replace_paragraph_text (self, paragraph, search_str, repl):
        """
            search all the runs in a paragraph and replace the search_str with repl
        """
        if search_str in paragraph.text:
            t_parags = len(paragraph.runs)
            for run_idx in range(t_parags):
                run = paragraph.runs[run_idx]
                if run.text.count('{')==run.text.count('}') and search_str in run.text:
                    run.text = run.text.replace(str(search_str), str(repl))
                elif run.text.count('{')!=run.text.count('}'):
                    #have a partial tag, need to merge runs
                    base_run = run
                    close_found = False
                    if run_idx < t_parags:
                        for mrg_idx in range(run_idx+1,t_parags):
                            m_run = paragraph.runs[mrg_idx]
                            base_run.text = base_run.text + m_run.text
                            if '}' in m_run.text:
                                close_found = True
                                break
                        if close_found:
                            #delete all extra runs
                            for i in reversed(range(run_idx+1,mrg_idx+1)):
                               self.delete_run(paragraph.runs[i]) 
                            self._replace_paragraph_text(paragraph, search_str, repl)
                            break

    def _merge_runs(self, paragraph):
        cur_text=''
        first=True
        for run in paragraph.runs:
            cur_text = cur_text + run.text
            if first != True:
                self.delete_run(run)
            first=False
        if len(paragraph.runs) > 0:   
            run = paragraph.runs[0]
            run.text = cur_text
        else: 
            run = paragraph.add_run()
        return run

    def update_chart(self, name:str,df:DataFrame):
        shape = self.get_shape_by_name(name)
        if shape != None:
            titles = list(df.index.values)
            data=df[df.columns[0]].values.tolist()
            for i in range(0,len(data)):
                if (isinstance(data[i],str)):
                    data[i] = int(data[i].replace(',',''))

            if shape.has_chart:
                chart_data = CategoryChartData()
#                chart_data = ChartData()
                chart_data.categories = titles
                chart_data.add_series('Series 1', tuple(data))
                shape.chart.replace_data(chart_data)
        else:
            self.log.warning(f'chart {name} not found in template')

    def replace_block(self, begin_tag, end_tag, repl_text):
        for slide in self._prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text
                        if text.find(begin_tag)!=-1:
                            run=self.merge_runs(paragraph)
                            run_text = run.text
                            text_prefix = text[:run_text.find(begin_tag)]
                            text_suffix = text[run_text.find(end_tag)+len(end_tag):]
                            new_text = text_prefix + repl_text + text_suffix
                            run.text = run.text.replace(run_text,new_text)

    from typing import Union

    def estimate_text_box_size(self,
            txt,
            font,  # ImageFont
            max_width: Union[int, None] = None,
            line_spacing: int = 4
    ):
        """
        Example of use:
        right_margin = left_margin = Length(Cm(0.25)).pt * pt_per_pixel
        top_margin = bottom_margin = Length(Cm(0.13)).pt * pt_per_pixel
        width, height = estimate_text_box_size(
            txt,
            font,
            max_width=shape_width - (right_margin + left_margin),
        )

        print("Computed in pixels (w, h)")
        print((width + right_margin + left_margin, height + top_margin + bottom_margin))


        :param txt:
        :param font:
        :param max_width:
        :param line_spacing:
        :return:
        """

        from PIL import ImageDraw, Image

        image = Image.new(size=(400, 300), mode='RGB')
        draw = ImageDraw.Draw(image)
        emu_per_inch = 914400
        px_per_inch = 72.0
        pt_per_pixel = 0.75

        fontsize_pt = 12
        # font = ImageFont.truetype("arial.ttf", int(fontsize_pt / pt_per_pixel))
        import textwrap, math
        if max_width is not None:
            actual_txt = []
            for line in txt.split("\n"):
                _, _, width, h = font.getbbox(line)
                split_at = len(line) // math.ceil(width / max_width)
                actual_txt = actual_txt + textwrap.wrap(line, width=split_at)

            new_lines = len(actual_txt)
            actual_txt = "\n".join(
                actual_txt
            )
        else:
            actual_txt = txt
            new_lines = 0

        left, top, right, bottom = draw.multiline_textbbox(
            (0, 0), actual_txt, font=font, spacing=line_spacing
        )
        ascent, descent = font.getmetrics()

        return right - left, bottom  # + descent * new_lines


    """ ***********************************************************************************************************
                                                Update Table Functionality
    *********************************************************************************************************** """
    def table_max_rows(self,tbl) -> int:
        #a notation of :## at the end of the table name indicates the max row count
        max_rows=-1
        splt = tbl.name.split(':')
        if len(splt) > 1:
            max_rows=int(splt[1])
            if max_rows == 0:
                max_rows=-1
                self.log.warning(f'Invalid max rows notation: {tbl.name}')
        return max_rows

    def update_table(self, name, df:DataFrame, include_index=True, 
                     background=None, forground=None,neg_num_bkg:list=None, 
                     header_rows=1,max_rows=-1):
        table_shape = self.get_shape_by_name(name)
        if table_shape is None:
            raise ValueError(f'Table not found in template: {name}')
        
        if max_rows < 0:
            max_rows = self.table_max_rows(table_shape)
        if max_rows < 0:
            max_rows = len(df)


        if table_shape != None and table_shape.has_table:
            table=table_shape.table

            colnames = list(df.columns)
            self.log.debug(f'filling table {name} with {len(df.index)} rows of data')

            # are there enough rows 
            rows, cols = df.shape

            trows = len(table._tbl.tr_lst)
            trows=trows+header_rows
            drows = len(df.index)
            if max_rows > 0:
                drows = min(rows,max_rows)
            
            """ calculate the number of rows in the table
                    * don't add more rows than the table reference max row count (table_name:max_rows)
                    * determin if the text will wrap, if so count as multiple rows
                        - the pptx packages does not directly support this
                        - perform a rough calculation using the font size and table cell margins
            """
            tbl_row_cnt = min(len(df),max_rows)+header_rows
            existing_rows = len(table._tbl.tr_lst)
            nrc = tbl_row_cnt-existing_rows
            #nrc = drows-trows+header_rows
            for r in range(nrc):
                self.add_row(table)

            # if trows < tbl_row_cnt:
            #     m = df.values
            #     existing_rows = len(table._tbl.tr_lst)
            #     nrc = tbl_row_cnt-existing_rows
            #     for row in range(rows):
            #         for col in range(cols):
            #             val = m[row, col]
            #             text = str(val) 

            #             # cell = table.cell(row+header_rows,col)
            #             # font = cell.text_frame.paragraphs[0].runs[0].font
            #             # width, height = self.estimate_text_box_size(text,font)                   
            #             pass


            # Insert the row zero names
            if include_index:
                for col_index, col_name in enumerate(df.index):
                    try:
                        # if has_header:
                        #     cell = table.cell(col_index+1,0)
                        # else:
                        cell = table.cell(col_index+header_rows,0)
                        text = str(col_name)
                        run = self._merge_runs(cell.text_frame.paragraphs[0])
                        run.text = text
                    except IndexError:
                        self.log.debug(f'index error in update_table ({name} @ row {col_index} with {text}) while setting df index')

            if background:
                cols = cols-1
            if forground:
                cols = cols-1

            if background:
                self.set_table_bgcolor(table,df[background],nrc,cols,header_rows)
            if forground:
                try:
                    self.set_table_font_color(table,df[forground],rows,cols,header_rows)
                except (KeyError):
                    self.log.warning(f'error setting forground for {name}')

            m = df.values
            for row in range(rows):
                for col in range(cols):
                    val = m[row, col]
                    text = str(val)
                    if include_index:
                        tbl_col=col+1
                    else:
                        tbl_col=col

                    try:
                        # if has_header:
                        #     cell = table.cell(row+1,tbl_col)
                        # else:
                        cell = table.cell(row+header_rows,tbl_col)

                        run = self._merge_runs(cell.text_frame.paragraphs[0])
                        run.text = text

                        test = 0
                        if type(val)==str and PowerPoint.isFloat(val.replace('%','')):
                            test = float(val.replace('%','').strip())

                        if (type(val)==str and PowerPoint.isFloat(test) and test < 0) or \
                           ((type(val)==int or type(val)==float) and int(val) < 0):
                            text_frame = cell.text_frame
                            run = self._merge_runs(text_frame.paragraphs[0])
                            run.font.color.rgb=RGBColor(255,255,255)
                            p = text_frame.paragraphs[0]

                            rPr = run._r.get_or_add_rPr()
                            hl = OxmlElement("a:highlight")
                            srgbClr = OxmlElement("a:srgbClr")

                            setattr(srgbClr,'val','FF0000')
                            hl.append(srgbClr)
                            rPr.append(hl)


                            pass

                    except IndexError:
                        self.log.debug(f'index error in update_table ({name}) while setting values')
                        break
    
    def set_table_bgcolor(self,table,colors,rows,cols,header_rows):
        for row in range(rows):
            rgb = colors.iloc[row].split(",")
            for col in range(cols):
                try:
                    # if has_header:
                    #     cell = table.cell(row+1,col)
                    # else:
                    cell = table.cell(row+header_rows,col)
                    cell.fill.solid() 
                    cell.fill.fore_color.rgb = RGBColor(int(rgb[0]), int(rgb [1]), int(rgb[2]))
                except IndexError:
                    self.log.warning('index error in update_table while setting background color')
    
    def set_table_font_color(self,table,colors,rows,cols,header_rows):
        for row in range(rows):
            rgb = colors.iloc[row].split(",")
            for col in range(cols):
                try:
                    # if has_header:
                    #     cell = table.cell(row+1,col)
                    # else:
                    cell = table.cell(row+header_rows,col)

                    paragraph = cell.text_frame.paragraphs[0]
                    run = self._merge_runs(paragraph)
                    run.font.color.rgb=RGBColor(int(rgb[0]), int(rgb [1]), int(rgb[2]))
                except IndexError:
                    self.log.warning('index error in update_table while setting background color')

    def change_paragraph_color(self,paragraph,rgb) -> None:
        run = self._merge_runs(paragraph)
        run.font.color.rgb=RGBColor(int(rgb[0]), int(rgb [1]), int(rgb[2]))

    def fill_text_box_color(self,name:str,color:RGBColor) -> None:
        shapes = PowerPoint.ppt.get_shape_by_name(name,all=True)
        if shapes is not None:
            for shape in shapes:
                if shape.shape_type in [ MSO_SHAPE_TYPE.TEXT_BOX, MSO_SHAPE_TYPE.AUTO_SHAPE]:
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = color
                else:
                    self.log.warning(f'Error filling background color for {name}')


    def isFloat(str):
        rslt = True
        try:
            float(str)
        except ValueError:
            rslt = False
        return rslt 

    def add_row(self,table: Table) -> _Row:
        new_row = deepcopy(table._tbl.tr_lst[-1]) 
        # duplicating last row of the table as a new row to be added

        for tc in new_row.tc_lst:
            cell = _Cell(tc, new_row.tc_lst)
            run = self._merge_runs(cell.text_frame.paragraphs[0])
            run.text = 'xxx' # defaulting cell contents to empty text

        table._tbl.append(new_row) 
        return new_row

    # def copy_slide(self,index=-1,template=None):
    #     if index<0 and template is None:
    #         raise KeyError('invalid parameters: either index or template are required')
    #     if template is not None:
    #         source = template
    #     else:
    #         source = self._prs.slides[index]

    #     # Append slide with source's layout. Then delete shapes to get a blank slide
    #     dest = self._prs.slides.add_slide(source.slide_layout)
    #     for shp in dest.shapes:
    #         shp.element.getparent().remove(shp.element)dest
    #     # Copy shapes from source, in order
    #     for shape in source.shapes:
    #         new_shape = deepcopy(shape.element)
    #         dest.shapes._spTree.insert_element_before(new_shape, 'p:extLst')
    #     # Copy rels from source
    #     for key, val in source.part.rels.items():
    #         target = val._target
    #         dest.part.rels.add_relationship(val.reltype, target, val.rId, val.is_external)

    #     # # Move appended slide into target_index
    #     # self._prs.slides.element.insert(target_index, prs.slides.element[-1])
    #     return dest


    def move_slide(self,old_index, new_index):
        xml_slides = self._prs.slides._sldIdLst
        slides = list(xml_slides)
        xml_slides.remove(slides[old_index])
        xml_slides.insert(new_index, slides[old_index])

    def copy_block(self, tag, prefix, count,slide=None):
        search_start = f'{{{tag}}}'
        search_end = f'{{end_{tag}}}'

        block = []

        slides = []
        if slide is not None:
            slides = [slide] 
        else:
            slides = self._prs.slides


        found=False
        for slide in slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text=paragraph.text
                        if not found and text.find(search_start)!=-1:
                            #is the end in the same paragraph?
                            if text.find(search_end)!=-1:
                                run=self._merge_runs(paragraph)
                                old_text = run.text
                                text_prefix = text[:old_text.find(search_end)]
                                text_suffix = text[old_text.find(search_end):]

                                sub_text = ""
                                for app_no in range (2,count+1):
                                    temp = old_text[old_text.find(search_start)+len(search_start):old_text.find(search_end)]
                                    for p in prefix:
                                        temp = temp.replace(f'{p}1',f'{p}{app_no}')
                                    sub_text = sub_text + ", " + temp
                                
                                new_text = text_prefix.replace(search_start,'') + sub_text + text_suffix.replace(search_end,'')
                                #new_text = text_prefix + sub_text + text_suffix

                                run.text = run.text.replace(old_text,new_text)
                            else:
                                found=True
                        elif found and text.find(search_end)!=-1:
                            found=False
                            for app_no in range(2,count+1):
                                self.paste_block(block, shape.text_frame,app_no)
                            if paragraph.text==search_end:
                                self.delete_paragraph(paragraph)
                            block=[]
                        if found:
                            if paragraph.text==search_start:
                                self.delete_paragraph(paragraph)
                            else:
                                block.append(paragraph)

    def delete_paragraph(self,paragraph):
        p = paragraph._p
        parent_element = p.getparent()
        parent_element.remove(p)

    def paste_block(self,block, text_frame,app_no):
        start_at = block[-1]

        for b in block:
            p = text_frame.add_paragraph()
            p.alignment = b.alignment
            p.line_spacing = b.line_spacing
            p.level = b.level

            for r in b.runs:    
                run = p.add_run()
                run.text = deepcopy(r.text)
                font = run.font
                font.name = r.font.name
                font.size = r.font.size
                font.bold = r.font.bold
                font.italic = r.font.italic
                if hasattr(r.font.color, 'rgb'):
                    font.color.rgb = r.font.color.rgb
                # else:
                #     font.color.theme_color = r.font.color.theme_color 

            run = self.merge_runs(p)
            run.text = run.text.replace("{app1_",f'{{app{app_no}_')
            run.text = run.text.replace("{end_app1_",f'{{end_app{app_no}_')

    def merge_runs(self, paragraph):
        cur_text=''
        first=True
        for run in paragraph.runs:
            cur_text = cur_text + run.text
            if first != True:
                self.delete_run(run)
            first=False
        if len(paragraph.runs) > 0:   
            run = paragraph.runs[0]
            run.text = cur_text
        else: 
            run = paragraph.add_run()
        return run

    def replace_shape_name (self, slide, search_str, repl_str):
        def update_name(shape,search_str,repl_str) -> None:
            shape.name = shape.name.replace(search_str,repl_str)

        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                for grp_shape in shape.shapes:
                    if grp_shape.name.find(search_str) != -1: 
                        update_name(grp_shape,search_str,repl_str)
            else:
                if shape.name.find(search_str) != -1: 
                    update_name(shape,search_str,repl_str)

    def get_page_no(self,shape):
        page_no = 0
        if shape:
            while True:
                shape = self.get_shape_parent(shape)
                if type(shape).__name__ == 'Slide':
                    break;
            
            page_no = self._prs.slides.index(shape) 
        return page_no

    def get_shape_parent(self,shape):
        rslt = None
        if hasattr(shape,'_parent'):
            rslt = shape._parent
        return rslt

    def get_slide(self,shape):
        while True:
            shape = self.ppt.get_shape_parent(shape) 
            if type(shape) is Slide:
                return shape
            elif shape is None:
                break
        return None

    def duplicate_slides(self, app_cnt):
        target_index = -1
        appendix=self.get_shape_by_name('StartAppendix')
        if appendix is not None:
            target_index = self._prs.slides.index(appendix._parent.parent)

        for cnt in range(2,app_cnt+1):
            for idx, slide in enumerate(self._prs.slides):
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            if paragraph.text == "{app_per_page}":
                                # new_slide = duplicate_slide(self._prs,slide)

                                new_slide = self.copy_slide(idx,target_index=target_index)
                                self._replace_slide_text(new_slide,"{app_per_page}","")
                                self._replace_slide_text(new_slide,"{app1_",f'{{app{cnt}_')
                                self.replace_shape_name(new_slide,"app1_",f'app{cnt}_')
                                target_index+=1

        for idx, slide in enumerate(self._prs.slides):
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        if paragraph.text == "{app_per_page}":
                            self._replace_slide_text(slide,"{app_per_page}","")

    def copy_slide(self,index=-1,template=None,target_index=-1):
        if index<0 and template is None:
            raise KeyError('invalid parameters: either index or template are required')

        if template is not None:
            source = template
        else:
            source = self._prs.slides[index]

        blank_slide_layout = source.slide_layout
        dest = self._prs.slides.add_slide(blank_slide_layout)
        if target_index >= 0:
            self.move_slide(self._prs.slides.index(dest),target_index)

        for shp in source.shapes:
            el = shp.element
            newel = deepcopy(el)
            dest.shapes._spTree.insert_element_before(newel, 'p:extLst')

        for key, value in source.part.rels.items():
            try:
                # Make sure we don't copy a notesSlide relation as that won't exist
                if "notesSlide" not in value.reltype:
                    target = value._target
                    # if the relationship was a chart, we need to duplicate the embedded chart part and xlsx
                    if "chart" in value.reltype:
                        partname = target.package.next_partname(ChartPart.partname_template)
                        xlsx_blob = target.chart_workbook.xlsx_part.blob
                        target = ChartPart(partname, target.content_type,
                                        deepcopy(target._element), package=target.package)
                        target.chart_workbook.xlsx_part = EmbeddedXlsxPart.new(xlsx_blob, target.package)

                    dest.part.rels.add_relationship(value.reltype,
                                                    target,
                                                    value.rId)

            except AttributeError as err:
                self.log.error(f'Attribute Error {err} while copying slide {index} part {key} ({err})')
            except KeyError as ex:
                self.log.error(f'KeyError {ex} while copying slide {index} part {key}')

        return dest

